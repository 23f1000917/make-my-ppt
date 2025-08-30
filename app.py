import io
import json
import os
import re
import time
import requests
from typing import Any, Dict, List, Optional
from fastapi import FastAPI, UploadFile, Form, HTTPException
from fastapi.responses import HTMLResponse, StreamingResponse, Response
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE_TYPE
from openai import OpenAI
import anthropic
from google import genai

app = FastAPI(title="MakeMyPPT - AI Presentation Generator")

# Configuration
MAX_TEXT_LENGTH = 60000
MIN_SLIDES = 10
MAX_SLIDES = 40
MAX_TEMPLATE_SIZE = 30 * 1024 * 1024
DEFAULT_MODELS = {
    "openai": "gpt-4o-mini",
    "aipipe": "gpt-4o-mini", 
    "anthropic": "claude-3-5-sonnet-latest",
    "gemini": "gemini-2.5-flash"
}

@app.get("/", response_class=HTMLResponse)
async def home():
    """Serve the main interface"""
    html_path = os.path.join(os.path.dirname(__file__), "index.html")
    if os.path.exists(html_path):
        with open(html_path, "r", encoding="utf-8") as f:
            return HTMLResponse(content=f.read())
    return HTMLResponse(content="<h1>MakeMyPPT</h1><p>Interface file not found</p>", status_code=404)


@app.post("/create")
async def create_presentation(
    content: str = Form(...),
    style_guide: Optional[str] = Form(None),
    ai_provider: str = Form(...),
    api_key: str = Form(...),
    model: Optional[str] = Form(None),
    slide_count: Optional[int] = Form(None),
    reuse_images: bool = Form(False),
    template: Optional[UploadFile] = None,
):
    # Input validation
    if not content or not content.strip():
        raise HTTPException(status_code=400, detail="Content is required")

    # Validate slide count
    target_slides = None
    if slide_count is not None:
        target_slides = max(1, min(MAX_SLIDES, slide_count))

    # Process template
    template_data = None
    if template and template.filename:
        if not template.filename.lower().endswith((".pptx", ".potx")):
            raise HTTPException(status_code=400, detail="Invalid template format")
        template_data = await template.read()
        if len(template_data) > MAX_TEMPLATE_SIZE:
            raise HTTPException(status_code=400, detail="Template too large")

    # Generate slide plan
    try:
        plan = await generate_slide_plan(
            content=content.strip()[:MAX_TEXT_LENGTH],
            style_guide=(style_guide or "").strip(),
            provider=ai_provider.lower(),
            api_key=api_key.strip(),
            model=(model or "").strip(),
            target_slides=target_slides,
        )
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"AI processing error: {e}")

    # Adjust slide count
    if target_slides:
        plan = adjust_slide_count(plan, target=target_slides, max_slides=MAX_SLIDES)
    else:
        plan = ensure_min_slides(plan, min_slides=MIN_SLIDES, max_slides=MAX_SLIDES)

    # Build presentation
    try:
        presentation = build_presentation(
            template_data=template_data,
            plan=plan,
            reuse_images=bool(reuse_images),
        )
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Presentation creation failed: {e}")

    headers = {
        "Content-Disposition": 'attachment; filename="presentation.pptx"',
        "Content-Type": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
    }
    return StreamingResponse(io.BytesIO(presentation), headers=headers)

def create_ai_prompt(content: str, style: str, target_slides: Optional[int]) -> str:
    slide_count = (
        f"Create exactly {target_slides} slides" if target_slides
        else f"Create between {MIN_SLIDES} and {MAX_SLIDES} slides"
    )
    
    return f"""
Transform this content into presentation slides. Return JSON only.

Requirements:
{slide_count}
- Titles: ≤80 characters
- Bullet points: 3-6 per slide, ≤120 chars each
- Follow this style: "{style}"

JSON format:
{json.dumps({"slides": [{"title": "Slide Title", "bullets": ["Point 1", "Point 2"]}]})}

Content:
{content}
""".strip()

async def generate_slide_plan(
    content: str,
    style_guide: str,
    provider: str,
    api_key: str,
    model: Optional[str],
    target_slides: Optional[int],
) -> Dict[str, Any]:
    model_name = model or DEFAULT_MODELS.get(provider)
    prompt = create_ai_prompt(content, style_guide, target_slides)
    
    if provider == "aipipe":
        return call_aipipe_api(api_key, model_name, prompt)
    
    messages = [{"role": "user", "content": prompt}]
    
    if provider == "openai":
        client = OpenAI(api_key=api_key)
        response = client.responses.create(
            model=model_name,
            input=messages,
            temperature=0.2,
        )
        content = extract_openai_content(response)
        return parse_json_response(content)
        
    elif provider == "anthropic":
        client = anthropic.Anthropic(api_key=api_key)
        response = client.messages.create(
            model=model_name,
            max_tokens=2048,
            system="Return valid JSON only",
            messages=messages,
            temperature=0.2,
        )
        text = "".join(block.text for block in response.content if block.type == "text")
        return parse_json_response(text)
        
    elif provider == "gemini":
        client = genai.Client(api_key=api_key)
        response = client.models.generate_content(
            model=model_name,
            contents=prompt,
            config={"response_mime_type": "application/json"},
        )
        return parse_json_response(response.text or "")
        
    else:
        raise HTTPException(status_code=400, detail="Unsupported AI provider")

def call_aipipe_api(api_key: str, model: str, prompt: str) -> Dict[str, Any]:
    """Call AI Pipe API"""
    url = "https://aipipe.org/openrouter/v1/chat/completions"
    headers = {
        "Content-Type": "application/json",
        "Authorization": f"Bearer {api_key}",
    }
    data = {
        "model": model,
        "messages": [{"role": "user", "content": prompt}],
        "max_tokens": 1000,
        "temperature": 0.7
    }
    
    response = requests.post(url, headers=headers, json=data)
    if response.status_code != 200:
        raise RuntimeError(f"API error: {response.status_code}")
    
    try:
        result = response.json()
        content = result["choices"][0]["message"]["content"]
        return parse_json_response(content)
    except Exception:
        return {"slides": [{"title": "Presentation", "bullets": [content]}]}

def extract_openai_content(response) -> str:
    """Extract text from OpenAI response"""
    if hasattr(response, 'output_text'):
        return response.output_text
    try:
        if response.output and response.output[0].content:
            return response.output[0].content[0].text
    except Exception:
        pass
    return json.dumps(response, default=str)

def parse_json_response(text: str) -> Dict[str, Any]:
    """Parse JSON from AI response"""
    text = (text or "").strip()
    if not text:
        return {"slides": []}
    try:
        return json.loads(text)
    except Exception:
        # Try to extract JSON from text
        match = re.search(r"\{.*\}", text, flags=re.DOTALL)
        if match:
            return json.loads(match.group())
        raise

def ensure_min_slides(plan: Dict[str, Any], min_slides: int, max_slides: int) -> Dict[str, Any]:
    """Ensure minimum slide count"""
    slides = []
    for slide in plan.get("slides", []):
        title = str(slide.get("title", "")).strip() or "Slide"
        bullets = [str(b).strip() for b in slide.get("bullets", []) if str(b).strip()]
        slides.append({"title": title, "bullets": bullets})
    
    # Split slides with too many bullets
    i = 0
    while len(slides) < min_slides and i < len(slides):
        slide = slides[i]
        if len(slide["bullets"]) > 3:
            extra = slide["bullets"][3:]
            slide["bullets"] = slide["bullets"][:3]
            while extra and len(slides) < min_slides:
                chunk = extra[:3]
                extra = extra[3:]
                slides.insert(i + 1, {"title": f"{slide['title']} (cont.)", "bullets": chunk})
                i += 1
        i += 1
    
    # Add placeholder slides if needed
    while len(slides) < min_slides:
        slides.append({"title": f"Slide {len(slides)+1}", "bullets": []})
    
    plan["slides"] = slides[:max_slides]
    return plan

def adjust_slide_count(plan: Dict[str, Any], target: int, max_slides: int) -> Dict[str, Any]:
    """Adjust to exact slide count"""
    target = max(1, min(max_slides, target))
    slides = plan.get("slides", [])
    
    # Normalize slides
    normalized = []
    for slide in slides:
        normalized.append({
            "title": (str(slide.get("title", "")) or "Slide").strip(),
            "bullets": [str(b).strip() for b in slide.get("bullets", []) if str(b).strip()],
        })
    
    # Adjust count
    if len(normalized) < target:
        plan["slides"] = ensure_min_slides({"slides": normalized}, target, max_slides)["slides"][:target]
    elif len(normalized) > target:
        # Try to merge continuation slides
        merged = []
        i = 0
        while i < len(normalized):
            current = normalized[i]
            if (i + 1 < len(normalized) and 
                normalized[i + 1]["title"].startswith(current["title"])):
                # Merge continuation
                next_slide = normalized[i + 1]
                current["bullets"].extend(next_slide["bullets"])
                current["bullets"] = current["bullets"][:8]
                i += 2
                merged.append(current)
            else:
                merged.append(current)
                i += 1
        plan["slides"] = merged[:target]
    
    return plan

def build_presentation(
    template_data: Optional[bytes],
    plan: Dict[str, Any],
    reuse_images: bool = False,
) -> bytes:
    """Build PowerPoint presentation"""
    prs = Presentation(io.BytesIO(template_data)) if template_data else Presentation()
    
    # Collect images from template if needed
    template_images = []
    if template_data and reuse_images:
        for slide in prs.slides:
            images = []
            for shape in slide.shapes:
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    try:
                        images.append({
                            "blob": shape.image.blob,
                            "left": int(shape.left),
                            "top": int(shape.top),
                            "width": int(shape.width),
                            "height": int(shape.height),
                        })
                    except Exception:
                        pass
            template_images.append(images)
    
    # Clear existing slides
    for slide_id in list(prs.slides._sldIdLst):
        rId = slide_id.rId
        prs.part.drop_rel(rId)
        prs.slides._sldIdLst.remove(slide_id)
    
    # Find appropriate layout
    layout_idx = find_slide_layout(prs)
    slide_width = int(prs.slide_width)
    slide_height = int(prs.slide_height)
    
    # Create slides
    for idx, slide_data in enumerate(plan.get("slides", [])):
        title = str(slide_data.get("title", "")).strip()[:120] or f"Slide {idx+1}"
        bullets = [str(b).strip() for b in slide_data.get("bullets", []) if str(b).strip()][:10]
        
        slide = prs.slides.add_slide(prs.slide_layouts[layout_idx])
        
        # Add images first (behind text)
        if template_images and idx < len(template_images):
            text_zones = find_text_zones(slide)
            for img in template_images[idx]:
                img_rect = create_rect(img["left"], img["top"], img["width"], img["height"])
                if overlaps_text(img_rect, text_zones):
                    safe_zone = find_safe_zone(slide_width, slide_height, text_zones)
                    img_rect = fit_in_box(img_rect, safe_zone)
                try:
                    slide.shapes.add_picture(
                        io.BytesIO(img["blob"]),
                        img_rect["left"], img_rect["top"],
                        width=img_rect["width"], height=img_rect["height"]
                    )
                except Exception:
                    pass
        
        # Add title
        if slide.shapes.title:
            slide.shapes.title.text = title
        else:
            textbox = slide.shapes.add_textbox(Inches(0.7), Inches(0.5), Inches(9), Inches(1))
            textbox.text_frame.text = title
            textbox.text_frame.paragraphs[0].font.size = Pt(32)
            textbox.text_frame.paragraphs[0].font.bold = True
        
        # Add content
        body = None
        for placeholder in slide.placeholders:
            try:
                if placeholder.is_placeholder and placeholder.placeholder_format.type != 1:
                    body = placeholder
                    break
            except Exception:
                continue
        
        if body:
            tf = body.text_frame
            tf.clear()
            if bullets:
                tf.paragraphs[0].text = bullets[0]
                for bullet in bullets[1:]:
                    p = tf.add_paragraph()
                    p.text = bullet
        else:
            textbox = slide.shapes.add_textbox(Inches(0.8), Inches(2.0), Inches(8.5), Inches(4.5))
            tf = textbox.text_frame
            if bullets:
                tf.text = bullets[0]
                for bullet in bullets[1:]:
                    p = tf.add_paragraph()
                    p.text = bullet
    
    # Save and return
    output = io.BytesIO()
    prs.save(output)
    return output.getvalue()

def find_slide_layout(prs: Presentation) -> int:
    """Find appropriate slide layout"""
    for i, layout in enumerate(prs.slide_layouts):
        has_title = has_body = False
        for ph in layout.placeholders:
            try:
                ph_type = ph.placeholder_format.type
                if ph_type == 1:
                    has_title = True
                elif ph_type in (2, 7):
                    has_body = True
            except Exception:
                continue
        if has_title and has_body:
            return i
    return 0

def find_text_zones(slide) -> List[Dict[str, int]]:
    """Find all text areas on slide"""
    zones = []
    for shape in slide.shapes:
        try:
            if getattr(shape, "has_text_frame", False) or getattr(shape, "is_placeholder", False):
                zones.append(create_rect(
                    int(shape.left), int(shape.top), 
                    int(shape.width), int(shape.height)
                ))
        except Exception:
            continue
    return zones

def create_rect(left: int, top: int, width: int, height: int) -> Dict[str, int]:
    return {"left": left, "top": top, "width": width, "height": height}

def overlaps_text(rect: Dict[str, int], zones: List[Dict[str, int]], threshold: float = 0.1) -> bool:
    """Check if rectangle overlaps text zones"""
    rect_area = rect["width"] * rect["height"]
    for zone in zones:
        overlap = calculate_overlap(rect, zone)
        if overlap / rect_area > threshold:
            return True
    return False

def calculate_overlap(a: Dict[str, int], b: Dict[str, int]) -> int:
    """Calculate overlap area between two rectangles"""
    ax1, ay1 = a["left"], a["top"]
    ax2, ay2 = ax1 + a["width"], ay1 + a["height"]
    bx1, by1 = b["left"], b["top"]
    bx2, by2 = bx1 + b["width"], by1 + b["height"]
    
    ix1 = max(ax1, bx1)
    iy1 = max(ay1, by1)
    ix2 = min(ax2, bx2)
    iy2 = min(ay2, by2)
    
    if ix2 < ix1 or iy2 < iy1:
        return 0
    return (ix2 - ix1) * (iy2 - iy1)

def find_safe_zone(slide_w: int, slide_h: int, text_zones: List[Dict[str, int]]) -> Dict[str, int]:
    """Find safe area for images that doesn't overlap text"""
    padding = 91440  # ~0.1 inch
    return create_rect(padding, padding, slide_w - 2 * padding, slide_h - 2 * padding)

def fit_in_box(rect: Dict[str, int], box: Dict[str, int]) -> Dict[str, int]:
    """Scale rectangle to fit within box while maintaining aspect ratio"""
    rect_w, rect_h = rect["width"], rect["height"]
    box_w, box_h = box["width"], box["height"]
    
    scale = min(box_w / rect_w, box_h / rect_h, 1.0)
    new_w = int(rect_w * scale)
    new_h = int(rect_h * scale)
    new_left = box["left"] + (box_w - new_w) // 2
    new_top = box["top"] + (box_h - new_h) // 2
    
    return create_rect(new_left, new_top, new_w, new_h)
